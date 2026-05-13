from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pathlib import Path
from typing import List, Optional, Dict, Any
from docx import Document
from pptx import Presentation
import requests
import uuid
import tempfile
import os
import subprocess
import json
import re
from datetime import datetime

APP_DIR = Path(__file__).parent
OUTPUT_DIR = APP_DIR / "output"
SKILLS_DIR = APP_DIR / "skills"
AGENTS_DIR = APP_DIR / "agents"
RUNS_DIR = APP_DIR / "agent_runs"
OUTPUT_DIR.mkdir(exist_ok=True)
RUNS_DIR.mkdir(exist_ok=True)
AGENTS_DIR.mkdir(exist_ok=True)

OLLAMA_BASE_URL = os.getenv("OLLAMA_BASE_URL", "http://127.0.0.1:11434")
DEFAULT_MODEL = os.getenv("OLLAMA_MODEL", "deepseek-r1:8b")

# CORS: set CORS_ALLOW_ALL=true for wildcard (e.g. LAN/dev), or
# ALLOWED_ORIGINS=http://host1:3000,http://host2:3000 for explicit list.
_CORS_ALLOW_ALL = os.getenv("CORS_ALLOW_ALL", "false").lower() == "true"
_ALLOWED_ORIGINS_ENV = os.getenv("ALLOWED_ORIGINS", "")
if _CORS_ALLOW_ALL:
    _allowed_origins: list = ["*"]
elif _ALLOWED_ORIGINS_ENV:
    _allowed_origins = [o.strip() for o in _ALLOWED_ORIGINS_ENV.split(",") if o.strip()]
else:
    _allowed_origins = [
        "http://localhost:3000", "http://127.0.0.1:3000",
        "http://localhost:3005", "http://127.0.0.1:3005",
    ]

app = FastAPI(title="DeepSeek Skill Studio + Agent Studio", version="1.3.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=_allowed_origins,
    # credentials cannot be used with wildcard origin
    allow_credentials=not _CORS_ALLOW_ALL,
    allow_methods=["*"],
    allow_headers=["*"],
)


def agent_file() -> Path:
    return AGENTS_DIR / "agents.json"


def load_skill(skill_name: str) -> str:
    skill_path = SKILLS_DIR / skill_name / "SKILL.md"
    if not skill_path.exists():
        raise HTTPException(status_code=404, detail=f"Skill not found: {skill_name}")
    return skill_path.read_text(encoding="utf-8")


def list_skills() -> List[str]:
    if not SKILLS_DIR.exists():
        return []
    return sorted([p.name for p in SKILLS_DIR.iterdir() if p.is_dir() and (p / "SKILL.md").exists()])


def load_agents() -> Dict[str, Any]:
    path = agent_file()
    if not path.exists():
        return {}
    return json.loads(path.read_text(encoding="utf-8"))


def save_agents(agents: Dict[str, Any]) -> None:
    path = agent_file()
    path.write_text(json.dumps(agents, indent=2), encoding="utf-8")


def slugify(value: str) -> str:
    slug = re.sub(r"[^a-zA-Z0-9]+", "_", value.lower()).strip("_")
    return slug[:60] or f"agent_{uuid.uuid4().hex[:8]}"


def get_agent(agent_id: str) -> Dict[str, Any]:
    agents = load_agents()
    if agent_id not in agents:
        raise HTTPException(status_code=404, detail=f"Agent not found: {agent_id}")
    return agents[agent_id]


def call_ollama(model: str, skill: str, prompt: str, context: str = "", agent_addendum: str = "") -> str:
    final_prompt = f"""
Skill instructions:
{skill}

Agent addendum:
{agent_addendum}

User request:
{prompt}

Context:
{context}
""".strip()

    chat_url = f"{OLLAMA_BASE_URL}/api/chat"
    generate_url = f"{OLLAMA_BASE_URL}/api/generate"
    try:
        response = requests.post(
            chat_url,
            json={
                "model": model,
                "messages": [
                    {"role": "system", "content": f"{skill}\n\n{agent_addendum}"},
                    {"role": "user", "content": f"{prompt}\n\nContext:\n{context}"},
                ],
                "stream": False,
            },
            timeout=900,
        )
        if response.status_code == 404:
            raise requests.HTTPError("/api/chat not available", response=response)
        response.raise_for_status()
        payload = response.json()
        return payload.get("message", {}).get("content", "")
    except Exception:
        try:
            response = requests.post(
                generate_url,
                json={"model": model, "prompt": final_prompt, "stream": False},
                timeout=900,
            )
            response.raise_for_status()
            return response.json().get("response", "")
        except Exception as exc:
            raise HTTPException(
                status_code=500,
                detail=f"Ollama call failed. Confirm Ollama is running and model is pulled. Error: {exc}",
            )


def call_ollama_chat(model: str, messages: List[Dict[str, str]]) -> str:
    chat_url = f"{OLLAMA_BASE_URL}/api/chat"
    generate_url = f"{OLLAMA_BASE_URL}/api/generate"
    try:
        response = requests.post(chat_url, json={"model": model, "messages": messages, "stream": False}, timeout=900)
        if response.status_code == 404:
            raise requests.HTTPError("/api/chat not available", response=response)
        response.raise_for_status()
        return response.json().get("message", {}).get("content", "")
    except Exception:
        prompt = "\n".join([f"{m.get('role','user').upper()}: {m.get('content','')}" for m in messages])
        try:
            response = requests.post(generate_url, json={"model": model, "prompt": prompt, "stream": False}, timeout=900)
            response.raise_for_status()
            return response.json().get("response", "")
        except Exception as exc:
            raise HTTPException(status_code=500, detail=f"Ollama chat failed: {exc}")


def create_docx(content: str, output_path: Path, title: str = "Generated Document"):
    doc = Document()
    doc.add_heading(title, 0)
    for raw in content.splitlines():
        line = raw.strip()
        if not line:
            continue
        if line.startswith("# "):
            doc.add_heading(line[2:].strip(), level=1)
        elif line.startswith("## "):
            doc.add_heading(line[3:].strip(), level=2)
        elif line.startswith("### "):
            doc.add_heading(line[4:].strip(), level=3)
        elif line.startswith("- "):
            doc.add_paragraph(line[2:].strip(), style="List Bullet")
        else:
            doc.add_paragraph(line)
    doc.save(output_path)


def create_pptx(content: str, output_path: Path):
    prs = Presentation()
    chunks = content.split("## Slide")
    if len(chunks) <= 1:
        chunks = content.split("# Slide")
    if len(chunks) <= 1:
        chunks = ["Overview\n" + content]
    for chunk in chunks:
        lines = [line.strip() for line in chunk.splitlines() if line.strip()]
        if not lines:
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = lines[0].replace("#", "").replace(":", "").strip()[:120]
        slide.shapes.title.text = title or "Slide"
        body = slide.placeholders[1]
        body.text = ""
        tf = body.text_frame
        for line in lines[1:6]:
            clean = line.replace("- ", "").strip()
            if not clean:
                continue
            p = tf.add_paragraph()
            p.text = clean[:300]
            p.level = 0
    prs.save(output_path)


def safe_read_file(path: Path, max_chars: int = 12000) -> str:
    try:
        if path.stat().st_size > 2_000_000:
            return f"[Skipped large file: {path.name}]"
        return path.read_text(encoding="utf-8", errors="ignore")[:max_chars]
    except Exception:
        return f"[Could not read file: {path.name}]"


def ingest_github_repo(repo_url: str, branch: Optional[str] = None, include_paths: str = "", exclude_paths: str = "") -> str:
    if not repo_url:
        return ""
    include_tokens = [x.strip() for x in include_paths.split(",") if x.strip()]
    exclude_tokens = [x.strip() for x in exclude_paths.split(",") if x.strip()]
    allowed_ext = {".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".md", ".yml", ".yaml", ".json", ".sql", ".xml", ".html", ".css", ".go", ".rs", ".sh", ".txt"}
    with tempfile.TemporaryDirectory() as tmp:
        cmd = ["git", "clone", "--depth", "1"]
        if branch:
            cmd.extend(["--branch", branch])
        cmd.extend([repo_url, tmp])
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
        parts = [f"--- GitHub Repo: {repo_url} ---"]
        total = 0
        for p in Path(tmp).rglob("*"):
            if not p.is_file() or ".git" in p.parts:
                continue
            rel = str(p.relative_to(tmp))
            if include_tokens and not any(tok in rel for tok in include_tokens):
                continue
            if exclude_tokens and any(tok in rel for tok in exclude_tokens):
                continue
            if p.suffix.lower() not in allowed_ext:
                continue
            text = f"\n--- File: {rel} ---\n{safe_read_file(p)}"
            parts.append(text)
            total += len(text)
            if total > 220000:
                parts.append("\n[Repo context truncated to stay within local model limits]")
                break
        return "\n".join(parts)


async def build_context(files: List[UploadFile], github_url: str, github_branch: str, include_paths: str, exclude_paths: str) -> str:
    context_parts = []
    for file in files:
        raw = await file.read()
        text = raw.decode("utf-8", errors="ignore")[:30000]
        context_parts.append(f"\n--- Uploaded File: {file.filename} ---\n{text}")
    if github_url.strip():
        context_parts.append(ingest_github_repo(github_url.strip(), github_branch.strip() or None, include_paths, exclude_paths))
    return "\n".join(context_parts)


def write_output(content: str, output_type: str, title: str = "Generated Output") -> Path:
    file_id = str(uuid.uuid4())
    if output_type == "docx":
        output_path = OUTPUT_DIR / f"{file_id}.docx"
        create_docx(content, output_path, title)
    elif output_type == "pptx":
        output_path = OUTPUT_DIR / f"{file_id}.pptx"
        create_pptx(content, output_path)
    else:
        output_path = OUTPUT_DIR / f"{file_id}.md"
        output_path.write_text(content, encoding="utf-8")
    return output_path


def log_agent_run(payload: Dict[str, Any]) -> Path:
    run_id = payload.get("run_id") or str(uuid.uuid4())
    path = RUNS_DIR / f"{run_id}.json"
    path.write_text(json.dumps(payload, indent=2), encoding="utf-8")
    return path


def installed_models() -> List[str]:
    try:
        tags = requests.get(f"{OLLAMA_BASE_URL}/api/tags", timeout=5).json()
        models = []
        for item in tags.get("models", []):
            name = item.get("name")
            if name:
                models.append(name)
        return sorted(models)
    except Exception:
        return []


def infer_agent_spec_from_prompt(prompt: str, model: str) -> Dict[str, Any]:
    skills = list_skills()
    fallback_skill = "document_writer" if "document_writer" in skills else (skills[0] if skills else "document_writer")
    output = "pptx" if any(x in prompt.lower() for x in ["ppt", "presentation", "slide", "deck"]) else "docx"
    if any(x in prompt.lower() for x in ["markdown", "md", "analysis only", "chat"]):
        output = "md"

    system = """
Create one JSON object for a local AI agent. Return only valid JSON, no markdown.
Fields: name, description, default_skill, default_output, system_addendum, allowed_tools.
Allowed default_output values: docx, pptx, md.
Pick default_skill from the provided skills only.
allowed_tools can include files, github, docx, pptx, markdown, chat.
""".strip()
    user = f"Available skills: {skills}\nAgent creation request: {prompt}"
    try:
        raw = call_ollama_chat(model, [{"role": "system", "content": system}, {"role": "user", "content": user}])
        match = re.search(r"\{.*\}", raw, re.DOTALL)
        if match:
            spec = json.loads(match.group(0))
        else:
            spec = json.loads(raw)
    except Exception:
        spec = {
            "name": prompt.strip().split("\n")[0][:60] or "Custom Agent",
            "description": f"Custom agent created from prompt: {prompt[:120]}",
            "default_skill": fallback_skill,
            "default_output": output,
            "system_addendum": f"Act as a focused local agent for this objective: {prompt}. Ask clarifying questions only when necessary. Produce practical, structured output.",
            "allowed_tools": ["files", "github", "docx", "pptx", "markdown", "chat"],
        }

    default_skill = spec.get("default_skill") if spec.get("default_skill") in skills else fallback_skill
    default_output = spec.get("default_output") if spec.get("default_output") in ["docx", "pptx", "md"] else output
    return {
        "name": str(spec.get("name") or "Custom Agent")[:80],
        "description": str(spec.get("description") or "On-demand custom local agent.")[:300],
        "default_skill": default_skill,
        "default_output": default_output,
        "system_addendum": str(spec.get("system_addendum") or f"Act as a focused agent for: {prompt}")[:2000],
        "allowed_tools": spec.get("allowed_tools") if isinstance(spec.get("allowed_tools"), list) else ["files", "github", "docx", "pptx", "markdown", "chat"],
        "created_by_prompt": prompt,
        "created_at": datetime.utcnow().isoformat() + "Z",
    }


@app.get("/health")
def health():
    models = installed_models()
    return {"status": "ok", "ollama_ok": bool(models), "ollama_base_url": OLLAMA_BASE_URL, "models": models}


@app.get("/models")
def models():
    return {"models": installed_models(), "default_model": DEFAULT_MODEL}


@app.get("/skills")
def skills():
    return {"skills": list_skills()}


@app.get("/agents")
def agents():
    return {"agents": load_agents()}


@app.post("/agent/create")
def create_agent(
    creation_prompt: str = Form(...),
    model: str = Form(DEFAULT_MODEL),
    agent_id: str = Form(""),
):
    spec = infer_agent_spec_from_prompt(creation_prompt, model)
    agents = load_agents()
    base_id = slugify(agent_id or spec["name"])
    new_id = base_id
    counter = 2
    while new_id in agents:
        new_id = f"{base_id}_{counter}"
        counter += 1
    agents[new_id] = spec
    save_agents(agents)
    return {"agent_id": new_id, "agent": spec, "agents": agents}


@app.post("/agent/run")
async def agent_run(
    prompt: str = Form(...),
    agent_id: str = Form(...),
    output_type: str = Form(""),
    model: str = Form(DEFAULT_MODEL),
    github_url: str = Form(""),
    github_branch: str = Form(""),
    include_paths: str = Form(""),
    exclude_paths: str = Form("node_modules,.git,dist,build,.venv,__pycache__"),
    files: List[UploadFile] = File(default=[]),
):
    agent = get_agent(agent_id)
    skill_name = agent.get("default_skill", "document_writer")
    selected_output = output_type or agent.get("default_output", "md")
    context = await build_context(files, github_url, github_branch, include_paths, exclude_paths)
    skill = load_skill(skill_name)
    result = call_ollama(model, skill, prompt, context, agent.get("system_addendum", ""))
    output_path = write_output(result, selected_output, agent.get("name", "Agent Output"))
    run_id = str(uuid.uuid4())
    log_agent_run({
        "run_id": run_id,
        "timestamp": datetime.utcnow().isoformat() + "Z",
        "agent_id": agent_id,
        "agent_name": agent.get("name"),
        "skill_name": skill_name,
        "model": model,
        "output_type": selected_output,
        "github_url": github_url,
        "include_paths": include_paths,
        "exclude_paths": exclude_paths,
        "download_file": output_path.name,
        "prompt": prompt,
        "raw_markdown": result,
    })
    return {"run_id": run_id, "filename": output_path.name, "download_url": f"/download/{output_path.name}", "raw_markdown": result}


@app.post("/chat")
async def chat(
    message: str = Form(...),
    model: str = Form(DEFAULT_MODEL),
    agent_id: str = Form(""),
    history_json: str = Form("[]"),
    github_url: str = Form(""),
    github_branch: str = Form(""),
    include_paths: str = Form(""),
    exclude_paths: str = Form("node_modules,.git,dist,build,.venv,__pycache__"),
    files: List[UploadFile] = File(default=[]),
):
    try:
        history = json.loads(history_json or "[]")
        if not isinstance(history, list):
            history = []
    except Exception:
        history = []

    agent_instruction = "You are a helpful local DeepSeek chat assistant. Be direct and practical."
    if agent_id:
        agent = get_agent(agent_id)
        skill = load_skill(agent.get("default_skill", "document_writer"))
        agent_instruction = f"{skill}\n\nAgent behavior:\n{agent.get('system_addendum', '')}"

    context = await build_context(files, github_url, github_branch, include_paths, exclude_paths)
    messages = [{"role": "system", "content": agent_instruction}]
    if context:
        messages.append({"role": "user", "content": f"Use this context for this chat session:\n{context}"})
    for h in history[-12:]:
        role = h.get("role") if h.get("role") in ["user", "assistant"] else "user"
        content = str(h.get("content", ""))[:12000]
        if content:
            messages.append({"role": role, "content": content})
    messages.append({"role": "user", "content": message})
    reply = call_ollama_chat(model, messages)
    return {"reply": reply}


@app.get("/agent/runs")
def agent_runs():
    runs = []
    for path in sorted(RUNS_DIR.glob("*.json"), key=lambda p: p.stat().st_mtime, reverse=True)[:50]:
        try:
            payload = json.loads(path.read_text(encoding="utf-8"))
            payload.pop("raw_markdown", None)
            runs.append(payload)
        except Exception:
            pass
    return {"runs": runs}


@app.post("/generate")
async def generate(
    prompt: str = Form(...),
    skill_name: str = Form(...),
    output_type: str = Form(...),
    model: str = Form(DEFAULT_MODEL),
    github_url: str = Form(""),
    github_branch: str = Form(""),
    include_paths: str = Form(""),
    exclude_paths: str = Form("node_modules,.git,dist,build,.venv,__pycache__"),
    files: List[UploadFile] = File(default=[]),
):
    try:
        context = await build_context(files, github_url, github_branch, include_paths, exclude_paths)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"GitHub or file ingestion failed: {exc}")
    skill = load_skill(skill_name)
    result = call_ollama(model, skill, prompt, context)
    output_path = write_output(result, output_type)
    return {"filename": output_path.name, "download_url": f"/download/{output_path.name}", "raw_markdown": result}


@app.get("/download/{filename}")
def download(filename: str):
    path = OUTPUT_DIR / filename
    if not path.exists():
        raise HTTPException(status_code=404, detail="File not found")
    return FileResponse(path, filename=filename)
